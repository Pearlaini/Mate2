# -*- coding: utf-8 -*-
import json
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

path = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_백업.pptx"
)
prs = Presentation(str(path))

def ovals(n):
    s = prs.slides[n-1]
    return sorted(
        [(sh.name, sh.text.strip()) for sh in s.shapes if "타원" in sh.name and sh.has_text_frame],
        key=lambda x: x[0],
    )

def title_line(n):
    for sh in prs.slides[n-1].shapes:
        if sh.has_text_frame and len(sh.text) > 50:
            return sh.text.split("\n")[0]
    return None

rows = []
for n in list(range(17, 27)) + [32, 33, 34]:
    rows.append({"slide": n, "title": title_line(n), "ovals": ovals(n)})

Path(__file__).resolve().parent.joinpath("_numbering_patterns.json").write_text(
    json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8"
)
