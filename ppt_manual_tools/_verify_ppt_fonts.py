# -*- coding: utf-8 -*-
"""수정본 슬라이드 서체 검증"""
import json
from pathlib import Path

from pptx import Presentation

PPT = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_수정.pptx"
)
OUT = Path(__file__).resolve().parent / "_ppt_font_check.json"

EXPECTED = {
    "title": ("Noto Sans KR Medium", 12.0),
    "subtitle": ("Noto Sans KR Medium", 10.0),
    "body": ("Noto Sans KR Light", 10.0),
}


def collect_runs(slide_no: int, prs) -> list[dict]:
    rows = []
    for sh in prs.slides[slide_no - 1].shapes:
        if not sh.has_text_frame:
            continue
        if "상품전시 관리" in sh.text and len(sh.text) < 20:
            continue
        for pi, p in enumerate(sh.text_frame.paragraphs):
            for ri, r in enumerate(p.runs):
                if not r.text:
                    continue
                rows.append(
                    {
                        "shape": sh.name,
                        "para": pi,
                        "run": ri,
                        "text": r.text[:60],
                        "font": r.font.name,
                        "size": r.font.size.pt if r.font.size else None,
                    }
                )
    return rows


def main() -> None:
    prs = Presentation(str(PPT))
    data = {str(n): collect_runs(n, prs) for n in [20, 21, 23, 41]}
    bad = []
    for sn, runs in data.items():
        for r in runs:
            if r["font"] not in ("Noto Sans KR Medium", "Noto Sans KR Light"):
                bad.append({**r, "slide": sn, "issue": "unknown_font"})
            if r["size"] not in (10.0, 12.0, 8.0):
                bad.append({**r, "slide": sn, "issue": "unexpected_size"})
    OUT.write_text(
        json.dumps({"runs": data, "issues": bad}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"issues: {len(bad)} -> {OUT}")


if __name__ == "__main__":
    main()
