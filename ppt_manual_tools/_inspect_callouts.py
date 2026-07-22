# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation

MANUAL_DIR = Path(r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0")
backup = sorted(MANUAL_DIR.glob("*교육ing_백업.pptx"))[-1]
modified = sorted(MANUAL_DIR.glob("*교육ing_수정.pptx"))[-1]

for label, path in [("backup", backup), ("modified", modified)]:
    prs = Presentation(str(path))
    for sn in [19, 20, 21, 22]:
        ovals = sorted(
            [
                (sh.name, sh.text.strip())
                for sh in prs.slides[sn - 1].shapes
                if "타원" in sh.name and sh.has_text_frame
            ],
            key=lambda x: x[0],
        )
        print(label, sn, ovals)

# slide 20 first lines
prs = Presentation(str(modified))
for sh in prs.slides[19].shapes:
    if sh.has_text_frame and "지급" in sh.text:
        print("slide20 lines:", sh.text.split("\n")[:12])
