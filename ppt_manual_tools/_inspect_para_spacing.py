# -*- coding: utf-8 -*-
import json
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

BACKUP = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_백업.pptx"
)
MODIFIED = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_수정.pptx"
)
ALT = MODIFIED.with_name(MODIFIED.stem + "_간격수정.pptx")
OUT = Path(__file__).resolve().parent / "_para_spacing.json"


def rows(path, slide_no, keyword):
    prs = Presentation(str(path))
    for sh in prs.slides[slide_no - 1].shapes:
        if sh.has_text_frame and keyword in sh.text:
            result = []
            for i, p in enumerate(sh.text_frame.paragraphs):
                ppr = p._p.find(qn("a:pPr"))
                ln = ppr.find(qn("a:lnSpc")) if ppr is not None else None
                aft = ppr.find(qn("a:spcAft")) if ppr is not None else None
                result.append(
                    {
                        "i": i,
                        "text": p.text[:55],
                        "has_pPr": ppr is not None,
                        "marL": ppr.get("marL") if ppr is not None else None,
                        "lnSpc": ln is not None,
                        "spcAft": aft is not None,
                    }
                )
            return result
    return []


candidates = [p for p in (MODIFIED, ALT) if p.exists()]
modified_path = max(candidates, key=lambda p: p.stat().st_mtime) if candidates else MODIFIED
data = {
    "backup_20": rows(BACKUP, 20, "지급 정보"),
    "modified_20": rows(modified_path, 20, "지급 정보"),
    "modified_41": rows(modified_path, 41, "사은품 주문서"),
    "modified_path": str(modified_path),
}
OUT.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
missing = sum(1 for r in data["modified_20"] if not r["has_pPr"])
print(f"missing pPr on slide20: {missing}/{len(data['modified_20'])}")
