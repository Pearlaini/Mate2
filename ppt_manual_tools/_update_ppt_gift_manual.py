# -*- coding: utf-8 -*-
"""사은품 지급 설정 매뉴얼 PPT 부분 수정 (톤&매너·번호 체계 유지)"""
import re
from pathlib import Path

from pptx import Presentation

from _ppt_styles import rebuild_shape_styled, remap_runs_in_shape

MANUAL_DIR = Path(r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0")


def _resolve_ppt_paths() -> tuple[Path, Path]:
    backups = sorted(MANUAL_DIR.glob("*교육ing_백업.pptx"))
    outs = sorted(MANUAL_DIR.glob("*교육ing_수정.pptx"))
    if not backups:
        raise FileNotFoundError("백업 PPT(*교육ing_백업.pptx)를 찾지 못했습니다.")
    backup = backups[-1]
    out = outs[-1] if outs else backup.with_name(backup.stem.replace("_백업", "_수정") + ".pptx")
    return backup, out


BACKUP, OUT = _resolve_ppt_paths()

# 19페이지 콜아웃이 9번까지 — 20페이지부터 이어서 10번부터
CALLOUT_CONTINUE_AFTER = 9

SLIDE20_INSERT_LINES = [
    "ⓔ 자동 적용 : ‘예’ 또는 ‘아니오’ 중 하나를 선택합니다.",
    "ⓕ 사은품 적용 유형 : 드롭다운에서 적용 대상을 선택합니다.",
    "① 주문 전체 적용 : 모든 주문서에 적용합니다.",
    "② 클레임없는 주문서만 적용 : 클레임이 없는 주문서만 적용합니다.",
    "③ 교환클레임 주문서만 적용 : 교환 클레임 주문서만 적용합니다.",
]

SLIDE21_REMAP = [
    ("ⓕ", "ⓖ"),
    ("ⓖ", "ⓗ"),
    ("ⓗ", "ⓘ"),
    ("ⓘ", "ⓙ"),
    ("ⓙ", "ⓚ"),
]

SLIDE23_INSERT_LINES = [
    "ⓔ 자동 적용 : ‘예’ 또는 ‘아니오’ 중 하나를 선택합니다.",
    "ⓕ 사은품 적용 유형 : 드롭다운에서 적용 대상을 선택합니다.",
]

PAGE41_NOTE = (
    "※ [상품 관리 > 사은품 지급 설정] 이벤트 등록 시 설정한 ‘사은품 적용 유형’"
    "(주문 전체 적용 / 클레임없는 주문서만 적용 / 교환클레임 주문서만 적용)에 따라 "
    "사은품 적용 대상 주문서가 결정됩니다."
)

SECTION2_TITLE = "2. 사은품 지급 이벤트 등록"
SECTION1_TITLE = "1. 화면 개요"


def circled_digit(n: int) -> str:
    if 1 <= n <= 20:
        return chr(0x2460 + n - 1)
    return str(n)


def _main_content_shape(slide, *keywords: str):
    for sh in slide.shapes:
        if not (hasattr(sh, "text") and sh.text.strip()):
            continue
        if all(k in sh.text for k in keywords):
            return sh
    return None


def _lines_with_insert(lines: list[str], after: str, insert: list[str]) -> list[str]:
    out: list[str] = []
    for line in lines:
        out.append(line)
        if line.strip().startswith(after):
            out.extend(insert)
    return out


def _normalize_section2_first_line(lines: list[str]) -> list[str]:
    if lines and lines[0].strip() in ("사은품 지급 이벤트 등록", SECTION2_TITLE):
        lines[0] = SECTION2_TITLE
    return lines


def update_slide19_section_title(slide) -> None:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text.split("\n")[0].strip() in ("화면 개요", SECTION1_TITLE):
            lines = sh.text.split("\n")
            lines[0] = SECTION1_TITLE
            rebuild_shape_styled(sh, lines, treat_first_as_title=True)
            return


def apply_section2_title(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        first = sh.text.split("\n")[0].strip()
        if first in ("사은품 지급 이벤트 등록", SECTION2_TITLE):
            lines = _normalize_section2_first_line(sh.text.split("\n"))
            rebuild_shape_styled(sh, lines, treat_first_as_title=True)


def renumber_oval_callouts(slide, offset: int) -> None:
    """스크린샷 콜아웃(타원) 숫자를 이전 슬라이드에서 이어서 재부여"""
    for sh in slide.shapes:
        if "타원" not in sh.name or not sh.has_text_frame:
            continue
        t = sh.text.strip()
        if t.isdigit():
            new_n = int(t) + offset
            sh.text = str(new_n)


def update_circled_callout_refs(slide, offset: int) -> None:
    """본문 속 원문자 숫자(⑤번 등)를 콜아웃 번호와 맞춤"""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                def repl(m: re.Match) -> str:
                    old = ord(m.group(1)) - 0x2460 + 1
                    return f"{circled_digit(old + offset)}번"

                r.text = re.sub(r"([\u2460-\u2473])번", repl, r.text)


def update_slide20(slide) -> None:
    sh = _main_content_shape(slide, "지급 정보", "ⓓ 선착순")
    if not sh:
        raise RuntimeError("슬라이드 20: 대상 텍스트 박스를 찾지 못했습니다.")
    lines = sh.text.split("\n")
    new_lines = _lines_with_insert(lines, "ⓓ 선착순", SLIDE20_INSERT_LINES)
    new_lines = _normalize_section2_first_line(new_lines)
    rebuild_shape_styled(sh, new_lines, treat_first_as_title=True)


def update_slide21(slide) -> None:
    for sh in slide.shapes:
        if hasattr(sh, "text") and (
            "지급 조건" in sh.text or "선택 주문 상품 관리" in sh.text
        ):
            remap_runs_in_shape(sh, SLIDE21_REMAP)


def update_slide23(slide) -> None:
    sh = _main_content_shape(slide, "ⓓ 선착순", "ⓐ 이벤트")
    if not sh:
        for sh in slide.shapes:
            if hasattr(sh, "text") and "ⓓ 선착순" in sh.text:
                break
        else:
            raise RuntimeError("슬라이드 23: 대상 텍스트 박스를 찾지 못했습니다.")
    lines = sh.text.split("\n")
    new_lines = _lines_with_insert(lines, "ⓓ 선착순", SLIDE23_INSERT_LINES)
    rebuild_shape_styled(sh, new_lines, treat_first_as_title=False)


def update_slide41(slide) -> None:
    sh = _main_content_shape(slide, "사은품 주문서 추가", "사은품 적용 팝업")
    if not sh:
        raise RuntimeError("슬라이드 41: 대상 텍스트 박스를 찾지 못했습니다.")
    lines = sh.text.split("\n")
    new_lines: list[str] = []
    for line in lines:
        new_lines.append(line)
        if line.strip() == "사은품 적용 팝업창에서 사은품 적용 유형을 선택합니다":
            new_lines.append(PAGE41_NOTE)
    rebuild_shape_styled(sh, new_lines)


def apply_callout_numbering(prs) -> None:
    offset = CALLOUT_CONTINUE_AFTER
    for slide_idx in (19, 20, 21):
        renumber_oval_callouts(prs.slides[slide_idx], offset)
    update_circled_callout_refs(prs.slides[21], offset)


def main() -> None:
    global BACKUP, OUT
    BACKUP, OUT = _resolve_ppt_paths()
    if not BACKUP.exists():
        raise FileNotFoundError(BACKUP)

    prs = Presentation(str(BACKUP))
    update_slide19_section_title(prs.slides[18])
    update_slide20(prs.slides[19])
    update_slide21(prs.slides[20])
    apply_section2_title(prs.slides[20])
    apply_section2_title(prs.slides[21])
    update_slide23(prs.slides[22])
    update_slide41(prs.slides[40])
    apply_callout_numbering(prs)

    try:
        prs.save(str(OUT))
        print(f"수정본 저장: {OUT}")
    except PermissionError:
        alt = OUT.with_name(OUT.stem + "_간격수정.pptx")
        prs.save(str(alt))
        print(f"원본 수정본이 열려 있어 대체 경로에 저장: {alt}")


if __name__ == "__main__":
    main()
