# -*- coding: utf-8 -*-
"""PPT 슬라이드 텍스트 추출 (임시 스크립트)"""
import json
import sys
from pathlib import Path

from pptx import Presentation

PPT_PATH = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing.pptx"
)
OUT_PATH = Path(__file__).resolve().parent / "_ppt_extract.json"


def slide_text(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            lines.append(shape.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                row_txt = " | ".join(cell.text.strip() for cell in row.cells)
                if row_txt.strip():
                    lines.append("[TABLE] " + row_txt)
    return lines


def main() -> None:
    prs = Presentation(str(PPT_PATH))
    data = {"total": len(prs.slides), "slides": {}}
    targets = list(range(1, 6)) + list(range(17, 28)) + list(range(38, 43))
    for i in targets:
        if 1 <= i <= len(prs.slides):
            data["slides"][str(i)] = slide_text(prs.slides[i - 1])
    OUT_PATH.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"written: {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
