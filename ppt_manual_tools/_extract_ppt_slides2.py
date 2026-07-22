# -*- coding: utf-8 -*-
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_PATH = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing.pptx"
)
OUT = Path(__file__).resolve().parent / "_ppt_slide_detail.json"


def shape_info(shape, depth=0):
    info = {
        "type": str(shape.shape_type),
        "name": shape.name,
        "has_text": hasattr(shape, "text"),
    }
    if hasattr(shape, "text") and shape.text.strip():
        info["text"] = shape.text
        if shape.has_text_frame:
            paras = []
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    run_info = []
                    for r in p.runs:
                        run_info.append(
                            {
                                "text": r.text,
                                "font_name": r.font.name,
                                "size": r.font.size.pt if r.font.size else None,
                                "bold": r.font.bold,
                            }
                        )
                    paras.append({"text": p.text, "runs": run_info})
            info["paragraphs"] = paras
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            rows.append([c.text for c in row.cells])
        info["table"] = rows
    return info


prs = Presentation(str(PPT_PATH))
result = {}
for sn in [20, 23, 41]:
    slide = prs.slides[sn - 1]
    shapes = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        si = shape_info(sh)
        if "text" in si or "table" in si:
            shapes.append(si)
    result[str(sn)] = shapes

OUT.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
print("done", OUT)
