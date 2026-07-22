# -*- coding: utf-8 -*-
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

BACKUP = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_백업.pptx"
)
MODIFIED = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_수정.pptx"
)
OUT = Path(__file__).resolve().parent / "_slide19_20_numbers.json"


def shape_dump(sh):
    item = {"name": sh.name, "type": str(sh.shape_type)}
    if hasattr(sh, "text") and sh.text.strip():
        item["text"] = sh.text.strip()[:120]
    if sh.has_text_frame:
        paras = []
        for p in sh.text_frame.paragraphs:
            ppr = p._p.find(qn("a:pPr"))
            paras.append(
                {
                    "text": p.text[:80],
                    "marL": ppr.get("marL") if ppr is not None else None,
                    "indent": ppr.get("indent") if ppr is not None else None,
                }
            )
        if paras:
            item["paragraphs"] = paras
    return item


def slide_dump(prs, n):
    slide = prs.slides[n - 1]
    shapes = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        si = shape_dump(sh)
        if "text" in si or "paragraphs" in si:
            shapes.append(si)
    return shapes


def main():
    data = {}
    for label, path in [("backup", BACKUP), ("modified", MODIFIED)]:
        if not path.exists():
            continue
        prs = Presentation(str(path))
        data[f"{label}_19"] = slide_dump(prs, 19)
        data[f"{label}_20"] = slide_dump(prs, 20)
    OUT.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(OUT)


if __name__ == "__main__":
    main()
