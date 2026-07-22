# -*- coding: utf-8 -*-
import json
from pathlib import Path
from pptx import Presentation

OUT = Path(
    r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0\#Mate2.0매뉴얼_OMS_교육ing_수정.pptx"
)
prs = Presentation(str(OUT))
result = {}
for sn in [20, 21, 23, 41]:
    texts = []
    for sh in prs.slides[sn - 1].shapes:
        if hasattr(sh, "text") and sh.text.strip() and "상품전시 관리" not in sh.text:
            if len(sh.text) > 40 or sn == 41:
                texts.append(sh.text)
    result[str(sn)] = texts

(Path(__file__).resolve().parent / "_ppt_verify.json").write_text(
    json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
)
print("ok")
