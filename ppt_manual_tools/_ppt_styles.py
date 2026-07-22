# -*- coding: utf-8 -*-
"""Mate2.0 OMS 매뉴얼 PPT 톤&매너 — run 서체 + 문단 간격(pPr) 유지"""
import re
from copy import deepcopy
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

MANUAL_DIR = Path(r"c:\AINI\샵이지3.0\00.사용자매뉴얼\아워박스_#Mate2.0")


def _backup_ppt() -> Path:
    matches = sorted(MANUAL_DIR.glob("*교육ing_백업.pptx"))
    if not matches:
        raise FileNotFoundError("백업 PPT를 찾지 못했습니다.")
    return matches[-1]

FONT_TITLE = ("Noto Sans KR Medium", 12)
FONT_SUBTITLE = ("Noto Sans KR Medium", 10)
FONT_BODY = ("Noto Sans KR Light", 10)

_SUB_OPTION_PPR: Optional[object] = None


def style_run(run, kind: str) -> None:
    name, size = {
        "title": FONT_TITLE,
        "subtitle": FONT_SUBTITLE,
        "body": FONT_BODY,
    }[kind]
    run.font.name = name
    run.font.size = Pt(size)


def clear_paragraph_runs(paragraph) -> None:
    p_elm = paragraph._p
    for r in list(p_elm.findall(qn("a:r"))):
        p_elm.remove(r)


def set_paragraph_segments(paragraph, segments: Iterable[tuple[str, str]]) -> None:
    clear_paragraph_runs(paragraph)
    for text, kind in segments:
        if text == "":
            continue
        run = paragraph.add_run()
        run.text = text
        style_run(run, kind)


def _clone_ppr(paragraph) -> Optional[object]:
    ppr = paragraph._p.find(qn("a:pPr"))
    return deepcopy(ppr) if ppr is not None else None


def apply_ppr(paragraph, ppr_clone: Optional[object]) -> None:
    if ppr_clone is None:
        return
    p_elm = paragraph._p
    old = p_elm.find(qn("a:pPr"))
    if old is not None:
        p_elm.remove(old)
    p_elm.insert(0, deepcopy(ppr_clone))


def _is_circled_digit_line(line: str) -> bool:
    return bool(re.match(r"^[\u2460-\u2473]", line))


def detect_para_kind(
    line: str,
    prev_kind: Optional[str],
    line_index: int,
    *,
    treat_first_as_title: bool,
    prev_line: str = "",
) -> str:
    if line == "":
        if prev_kind in ("field", "sub_option"):
            return "blank_field"
        return "blank_block"

    if (
        prev_line
        and "> 영역" in prev_line
        and not line.startswith("<")
        and not re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ]", line)
        and " : " not in line
    ):
        return "body_detail"

    if line.startswith("※"):
        return "note"

    if line_index == 0 and treat_first_as_title and not re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ]", line):
        return "block"

    if line.startswith("<") and "영역" in line:
        return "block"

    if re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ] ", line) and " : " in line:
        return "field"

    if _is_circled_digit_line(line):
        return "sub_option"

    if " : " in line and not re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ]", line):
        if line.startswith(("주문 전체", "클레임", "교환클레임")):
            return "sub_option"
        if line.startswith(("사은품 적용", "선택 사은품 적용")):
            return "field"
        return "sub_option"

    return "block"


def capture_ppr_templates(shape, *, treat_first_as_title: bool = True) -> dict[str, object]:
    templates: dict[str, object] = {}
    prev_kind: Optional[str] = None
    prev_line = ""
    for i, p in enumerate(shape.text_frame.paragraphs):
        kind = detect_para_kind(
            p.text,
            prev_kind,
            i,
            treat_first_as_title=treat_first_as_title,
            prev_line=prev_line,
        )
        if kind not in templates:
            cloned = _clone_ppr(p)
            if cloned is not None:
                templates[kind] = cloned
        if p.text != "":
            prev_kind = kind
            prev_line = p.text
    return templates


def get_sub_option_ppr() -> Optional[object]:
    """슬라이드 21 하위 옵션 문단 간격 템플릿"""
    global _SUB_OPTION_PPR
    if _SUB_OPTION_PPR is not None:
        return _SUB_OPTION_PPR
    try:
        backup = _backup_ppt()
    except FileNotFoundError:
        return None
    prs = Presentation(str(backup))
    for sh in prs.slides[20].shapes:
        if sh.has_text_frame and "모든 상품이 일치" in sh.text:
            for p in sh.text_frame.paragraphs:
                if p.text.startswith("모든 상품이 일치"):
                    _SUB_OPTION_PPR = _clone_ppr(p)
                    return _SUB_OPTION_PPR
    return None


def _resolve_ppr(kind: str, templates: dict[str, object]) -> Optional[object]:
    if kind in templates:
        return templates[kind]
    if kind == "sub_option":
        sub = get_sub_option_ppr()
        if sub is not None:
            return sub
    if kind == "note" and "block" in templates:
        return templates["block"]
    return templates.get("block")


def _parse_field_line(line: str) -> list[tuple[str, str]]:
    m = re.match(r"^([ⓐ-⓯Ⓐ-Ⓩ]) (.+?) : (.+)$", line)
    if m:
        circ, label, desc = m.groups()
        return [(f"{circ} ", "body"), (label, "subtitle"), (" : ", "body"), (desc, "body")]
    return [(line, "body")]


def _parse_angle_section(line: str) -> list[tuple[str, str]]:
    m = re.match(r"^<(.+?)> 영역$", line)
    if m:
        return [("<", "body"), (m.group(1), "subtitle"), ("> 영역", "body")]
    return [(line, "body")]


def _parse_circled_sub_line(line: str) -> list[tuple[str, str]] | None:
    m = re.match(r"^([\u2460-\u2473]) (.+?) : (.+)$", line)
    if m:
        circ, label, desc = m.groups()
        return [(f"{circ} ", "body"), (label, "subtitle"), (" : ", "body"), (desc, "body")]
    return None


def _parse_option_line(line: str) -> list[tuple[str, str]]:
    if " : " in line:
        label, _, desc = line.partition(" : ")
        return [(label, "subtitle"), (" : ", "body"), (desc, "body")]
    return [(line, "body")]


def _parse_bracket_line(line: str) -> list[tuple[str, str]]:
    segments: list[tuple[str, str]] = []
    pos = 0
    for m in re.finditer(r"\[([^\]]+)\]", line):
        if m.start() > pos:
            segments.append((line[pos : m.start()], "body"))
        segments.extend([("[", "body"), (m.group(1), "subtitle"), ("]", "body")])
        pos = m.end()
    if pos < len(line):
        segments.append((line[pos:], "body"))
    return segments if segments else [(line, "body")]


def classify_line(
    line: str, line_index: int, *, treat_first_as_title: bool = True
) -> list[tuple[str, str]]:
    if line == "":
        return [("", "body")]
    if (
        line_index == 0
        and treat_first_as_title
        and not re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ]", line)
    ):
        return [(line, "title")]
    if line.startswith("<") and line.endswith("영역"):
        return _parse_angle_section(line)
    if re.match(r"^[ⓐ-⓯Ⓐ-Ⓩ] ", line) and " : " in line:
        return _parse_field_line(line)
    if _is_circled_digit_line(line):
        parsed = _parse_circled_sub_line(line)
        if parsed:
            return parsed
    if line.startswith("※"):
        return [(line, "body")]
    if "[" in line:
        return _parse_bracket_line(line)
    if " : " in line:
        return _parse_option_line(line)
    return [(line, "body")]


def rebuild_shape_styled(
    shape,
    lines: list[str],
    *,
    treat_first_as_title: bool = True,
    extra_templates: Optional[dict[str, object]] = None,
) -> None:
    """텍스트 박스를 서체·문단 간격 규칙에 맞게 다시 씁니다."""
    if not shape.has_text_frame:
        return

    templates = capture_ppr_templates(shape, treat_first_as_title=treat_first_as_title)
    if extra_templates:
        templates.update(extra_templates)

    tf = shape.text_frame
    tf.word_wrap = True

    while len(tf.paragraphs) > 1:
        el = tf.paragraphs[-1]._p
        el.getparent().remove(el)

    prev_kind: Optional[str] = None
    prev_line = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        kind = detect_para_kind(
            line,
            prev_kind,
            i,
            treat_first_as_title=treat_first_as_title,
            prev_line=prev_line,
        )
        apply_ppr(p, _resolve_ppr(kind, templates))
        set_paragraph_segments(
            p, classify_line(line, i, treat_first_as_title=treat_first_as_title)
        )
        if line != "":
            prev_kind = kind
            prev_line = line


def remap_runs_in_shape(shape, mapping: list[tuple[str, str]]) -> None:
    """run 텍스트만 치환 — 기존 서체·문단 간격 유지"""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            text = r.text
            for src, dst in mapping:
                text = text.replace(src, dst)
            r.text = text
